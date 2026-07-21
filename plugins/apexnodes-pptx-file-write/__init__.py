"""ApexNodes structured PowerPoint export tool (hc-326)."""

from __future__ import annotations

import os
import re
import uuid
from pathlib import Path
from typing import Any

from tools.registry import tool_error, tool_result

PPTX_FILE_WRITE_SCHEMA = {
    "name": "pptx_file_write",
    "description": (
        "Render structured slides to a .pptx file and return its local path. "
        "After calling this, put `MEDIA:<file_path>` on a separate line in your "
        "reply so the file is sent to the user. Slides support title, subtitle, "
        "bullet points, tables, notes, and local images."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "slides": {
                "type": "array",
                "description": (
                    "Slides. Each item may include title, subtitle, bullets, table "
                    "{columns, rows}, image_path, and notes."
                ),
            },
            "title": {
                "type": "string",
                "description": "Deck title / file name (defaults to Agent deck).",
            },
        },
        "required": ["slides"],
    },
}


def _check_pptx_file_write() -> bool:
    import importlib.util

    try:
        return importlib.util.find_spec("pptx") is not None
    except (ImportError, ValueError):
        return False


def _document_cache_dir() -> Path:
    home: str | None = None
    try:
        from hermes_constants import get_hermes_home

        home = str(get_hermes_home())
    except Exception:  # noqa: BLE001
        home = os.getenv("HERMES_HOME") or str(Path.home() / ".hermes")
    cache_dir = Path(home) / "cache" / "documents"
    cache_dir.mkdir(parents=True, exist_ok=True)
    return cache_dir


def _safe_filename(title: str) -> str:
    base = re.sub(r"[\\/\x00-\x1f]", " ", title or "").strip()
    base = re.sub(r"\s+", " ", base) or "Agent deck"
    return base[:60]


def _as_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    return str(value)


def _normalise_table(table: Any) -> tuple[list[str], list[list[str]]]:
    if not isinstance(table, dict):
        return [], []
    rows = table.get("rows") or []
    if not isinstance(rows, list):
        return [], []
    columns_raw = table.get("columns") or []
    columns = [_as_text(item.get("header") or item.get("key") if isinstance(item, dict) else item) for item in columns_raw]
    if not columns:
        keys: list[str] = []
        for row in rows:
            if isinstance(row, dict):
                for key in row:
                    text = str(key)
                    if text not in keys:
                        keys.append(text)
        columns = keys
    if not columns:
        max_len = max((len(row) for row in rows if isinstance(row, list | tuple)), default=0)
        columns = [f"Column {index + 1}" for index in range(max_len)]

    normalised_rows: list[list[str]] = []
    for row in rows:
        if isinstance(row, dict):
            normalised_rows.append([_as_text(row.get(column, "")) for column in columns])
        elif isinstance(row, list | tuple):
            normalised_rows.append([_as_text(row[index] if index < len(row) else "") for index in range(len(columns))])
    return columns, normalised_rows


def render_slides_to_pptx(slides: list[dict[str, Any]], title: str, dest: Path) -> Path:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.core_properties.title = title[:120]

    for index, slide_data in enumerate(slides):
        if not isinstance(slide_data, dict):
            raise ValueError("each slide must be an object")
        has_body = bool(slide_data.get("bullets") or slide_data.get("table") or slide_data.get("image_path"))
        layout = prs.slide_layouts[1] if has_body else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = _as_text(slide_data.get("title") or f"Slide {index + 1}")[:160]
        if not has_body and len(slide.placeholders) > 1:
            slide.placeholders[1].text = _as_text(slide_data.get("subtitle"))

        bullets = slide_data.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        if bullets and len(slide.placeholders) > 1:
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for bullet_index, bullet in enumerate(bullets):
                paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
                paragraph.text = _as_text(bullet)
                paragraph.level = 0
                paragraph.font.size = Pt(22)

        columns, rows = _normalise_table(slide_data.get("table"))
        if columns:
            top = Inches(3.1 if bullets else 1.7)
            left = Inches(0.7)
            width = Inches(8.6)
            height = Inches(0.5 + 0.35 * max(len(rows), 1))
            shape = slide.shapes.add_table(len(rows) + 1, len(columns), left, top, width, height)
            table = shape.table
            for column_index, column in enumerate(columns):
                table.cell(0, column_index).text = column
            for row_index, row in enumerate(rows, start=1):
                for column_index, value in enumerate(row):
                    table.cell(row_index, column_index).text = value

        image_path = _as_text(slide_data.get("image_path")).strip()
        if image_path:
            path = Path(image_path)
            if not path.exists():
                raise ValueError(f"image_path does not exist: {image_path}")
            slide.shapes.add_picture(str(path), Inches(6.4), Inches(1.5), width=Inches(3.0))

        notes = _as_text(slide_data.get("notes")).strip()
        if notes:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = notes

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT

    prs.save(str(dest))
    return dest


def _handle_pptx_file_write(args: dict, **_kwargs) -> str:
    if not isinstance(args, dict):
        return tool_error("pptx_file_write expects a JSON object argument")

    slides = args.get("slides")
    if not isinstance(slides, list) or not slides:
        return tool_error("slides must be a non-empty list")

    title = str(args.get("title") or "Agent deck").strip() or "Agent deck"
    try:
        cache_dir = _document_cache_dir()
        dest = cache_dir / f"{_safe_filename(title)}_{uuid.uuid4().hex[:8]}.pptx"
        render_slides_to_pptx(slides, title, dest)
    except ImportError:
        return tool_error("python-pptx is not installed in this runtime image")
    except Exception as exc:  # noqa: BLE001
        return tool_error(f"Failed to render the PowerPoint file: {exc}")

    path = str(dest)
    return tool_result(
        success=True,
        file_path=path,
        media_tag=f"MEDIA:{path}",
        instruction="在回复里单独一行输出 media_tag 的值即可把 PPT 文件发给用户;另附不超过五行的核心摘要。",
    )


def register(ctx):
    ctx.register_tool(
        name="pptx_file_write",
        toolset="doc_delivery",
        schema=PPTX_FILE_WRITE_SCHEMA,
        handler=_handle_pptx_file_write,
        check_fn=_check_pptx_file_write,
        requires_env=[],
        description=PPTX_FILE_WRITE_SCHEMA["description"],
        emoji="📽️",
    )
